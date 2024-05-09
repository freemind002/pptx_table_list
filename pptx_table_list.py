# -*- coding:utf-8 -*-
# 將事件編號、事件名稱、處理情形放入表格當中
import math
from typing import Any, Dict, List, Text, Union

import my_case_name_list
import polars as pl
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.dml import MSO_THEME_COLOR


class PTTXReport(object):
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_height, self.prs.slide_width = Cm(19.05), Cm(25.4)
        self.title_list = ["事件編號", "事件名稱", "處理結果"]

    def data_to_table(
        self,
        table,
        data_list: Union[List[Text], List[Dict[Text, Any]]],
        header: bool,
    ):
        """將資料寫入table中

        Args:
            table (_type_): 每一頁要操作的table
            data_list (Union[List[Text], List[Dict[Text, Any]]]): 有可能是header的資料或是事件的相關資料
            header (bool): 要操作的是header還是其他的內容
        """
        if header is True:
            for index, data in enumerate(data_list):
                cell = table.cell(0, index)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame
                para = tf.paragraphs[0]
                para.text = data
                para.font.size = Pt(16)
                para.font.name = "微軟正黑體"
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER  # 水平置中對齊
        else:
            for data_index, data in enumerate(data_list, 1):
                for value_index, (title, value) in enumerate(data.items()):
                    cell = table.cell(data_index, value_index)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf = cell.text_frame
                    para = tf.paragraphs[0]
                    para.text = value
                    para.font.size = Pt(16)
                    para.font.name = "微軟正黑體"
                    para.font.bold = False
                    para.alignment = PP_ALIGN.CENTER  # 水平置中對齊

    def run_all(self):
        n = int(input("請輸入？筆資料為一頁："))
        case_list = (
            pl.LazyFrame(
                {
                    self.title_list[1]: my_case_name_list.case_name_list,
                }
            )
            .with_row_index(self.title_list[0], offset=1)
            .with_columns(pl.lit("處理完成").alias(self.title_list[2]))
            .cast(pl.String)
            .collect()
            .to_dicts()
        )
        # print(case_list)
        # 總計需插入？頁，使用無條件進位處理
        case_num = len(case_list)
        print(case_num)
        blank_page = int(math.ceil(case_num / n))
        print(blank_page)
        for page in range(blank_page):
            # 用內置模板(0-10)添加一個全空的ppt頁面
            blank_slide_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_slide_layout)
            left, top, width, height = Cm(1.76), Cm(0.49), Cm(21.89), Cm(1.28)
            tBox = slide.shapes.add_textbox(left, top, width, height)
            tf = tBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "事件處理統計列表"
            font = run.font
            font.name = "微軟正黑體"
            font.size = Pt(39)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # 插入表格
            rows, cols, left, top, width, height = (
                n + 1,
                3,
                Cm(0.5),
                Cm(4.5),
                Cm(1),
                Cm(1.23),
            )
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            # 調整行高、列寬
            for index in range(rows):
                table.rows[index].height = (
                    Cm(12.6 / n * (1.1 / 1.26)) if index == 0 else Cm(12.6 / n)
                )
            table.columns[0].width, table.columns[1].width, table.columns[2].width = (
                Cm(2.8),
                Cm(13.6),
                Cm(7.86),
            )
            # 寫入表頭
            self.data_to_table(table, self.title_list, header=True)
            # 如果資料為n的倍數筆資料，處理方式
            if page + 1 < blank_page:
                self.data_to_table(
                    table, case_list[0 + page * n : (page + 1) * n], header=False
                )
            else:
                self.data_to_table(table, case_list[0 + page * n :], header=False)
        self.prs.save("test.pptx")
        print("PPTX製作完成")

    def main(self):
        current_except = None
        try:
            self.run_all()
        except Exception as e:
            print(e)
            current_except = e
        finally:
            if current_except:
                raise current_except


if __name__ == "__main__":
    PTTXReport().main()
